#!/usr/bin/env python
# -*- coding: utf-8 -*-
import pandas as pd
import matplotlib.pyplot as plt
from matplotlib.colors import LinearSegmentedColormap, Normalize
from matplotlib import colormaps
from matplotlib.colorbar import ColorbarBase
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
import numpy as np
import yaml

# Load configuration from config.yaml
def load_config(config_path="config/config.yaml"):
    with open(config_path, "r") as file:
        return yaml.safe_load(file)

# Load data from CSV or Excel
def load_data(input_file):
    if input_file.endswith(".csv"):
        return pd.read_csv(input_file)
    elif input_file.endswith(".xlsx"):
        return pd.read_excel(input_file)
    else:
        raise ValueError("Unsupported file format. Please use .csv or .xlsx.")

# Create colormap based on config
def create_colormap(config):
    gradient = config['color_scale']['gradient']
    if isinstance(gradient, str) and gradient in colormaps:
        # Use built-in colormap
        cmap = colormaps[gradient].resampled(100)
    else:
        # Create custom colormap from gradient
        colors = [color['color'] for color in gradient]
        cmap = LinearSegmentedColormap.from_list('custom_colormap', colors, N=100)
    return cmap

# Map data values to colors using a colormap
def map_color(value, min_val, max_val, cmap):
    if value is None or np.isnan(value):
        return None  # No color for missing values
    if max_val == min_val:
        norm_value = 0.5  # All values are the same, use midpoint color
    else:
        norm_value = (value - min_val) / (max_val - min_val)
        norm_value = min(max(norm_value, 0), 1)  # Clamp to [0, 1]
        norm_value = norm_value * 99.0 / 100.0 # Prevent the value from reaching exactly 1
    rgba = cmap(norm_value)
    r, g, b = [int(255 * x) for x in rgba[:3]]
    return RGBColor(r, g, b)

# Create a colorbar image and save it
def create_colorbar_image(output_image, config, orientation):
    # Define colormap
    cmap = create_colormap(config)

    # Normalize values
    norm = Normalize(vmin=config['color_scale']['min_value'], vmax=config['color_scale']['max_value'])

    # Create the colorbar plot
    if orientation == 'vertical':
        plt.figure(figsize=(1.2, 4))  # Tall, narrow figure
    elif orientation == 'horizontal':
        plt.figure(figsize=(5, 1.2))  # Wide, short figure

    ax = plt.gca()
    cbar = ColorbarBase(ax, cmap=cmap, norm=norm, orientation=orientation)
    cbar.ax.tick_params(labelsize=24)

    # Save as image
    plt.tight_layout()
    plt.savefig(output_image, bbox_inches='tight', dpi=300)
    plt.close()

# Add colorbar image to all slides in the presentation
def add_colorbar_to_all_slides(prs, image_path, left=7.5, top=1):
    for slide in prs.slides:
        slide.shapes.add_picture(image_path, Inches(left + 2.0), Inches(top + 2.0), height=Inches(2.0))  # Add image to slide

# Annotate PowerPoint slides with gene expression data
def annotate_ppt(expression_data, input_pptx, output_pptx, config):
    # Create colorbar image
    colorbar_image = "output/colorbar_vertical.png"
    create_colorbar_image(colorbar_image, config, 'vertical')
    colorbar_image = "output/colorbar_horizontal.png"
    create_colorbar_image(colorbar_image, config, 'horizontal')

    # Select coloarbar image
    colorbar_image = "output/colorbar_vertical.png"

    # Define colormap
    cmap = create_colormap(config)

    prs = Presentation(input_pptx)

    # Annotate slides with gene data
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text in expression_data:
                    expr_value = expression_data[text]
                    color = map_color(expr_value,
                                      config['color_scale']['min_value'],
                                      config['color_scale']['max_value'],
                                      cmap)
                    if color is not None:
                        fill = shape.fill
                        fill.solid()
                        fill.fore_color.rgb = color

    # Add colorbar to all slides
    add_colorbar_to_all_slides(prs, colorbar_image)

    prs.save(output_pptx)

# Main function
def main(input_file, input_pptx, output_pptx, config_path="config/config.yaml"):
    config = load_config(config_path)
    data = load_data(input_file)

    # Convert data to a dictionary for mapping
    expression_data = dict(zip(data['Gene'], data['Value']))

    annotate_ppt(expression_data, input_pptx, output_pptx, config)
    print(f"Annotated pathway saved to {output_pptx}")

if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(description="Annotate pathway diagrams with gene expression data.")
    parser.add_argument("--input", required=True, help="Input CSV or Excel file with gene data.")
    parser.add_argument("--pathway", required=True, help="Input PowerPoint file with pathway diagram.")
    parser.add_argument("--output", required=True, help="Output PowerPoint file.")
    parser.add_argument("--config", default="config/config.yaml", help="Path to configuration file.")

    args = parser.parse_args()

    main(args.input, args.pathway, args.output, args.config)
